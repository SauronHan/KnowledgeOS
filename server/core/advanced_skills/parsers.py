import os
from pathlib import Path

def parse_pdf(file_path: Path) -> str:
    from pypdf import PdfReader
    text = []
    reader = PdfReader(str(file_path))
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)
    return "\n\n".join(text)

def parse_docx(file_path: Path) -> str:
    import docx
    doc = docx.Document(str(file_path))
    text = []
    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text.strip())
    return "\n\n".join(text)

def parse_pptx(file_path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(file_path))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n\n".join(text)

def parse_xlsx(file_path: Path) -> str:
    import pandas as pd
    excel_file = pd.ExcelFile(str(file_path))
    text = []
    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(str(file_path), sheet_name=sheet_name)
        text.append(f"### Sheet: {sheet_name}")
        text.append(df.to_markdown(index=False))
    return "\n\n".join(text)

def parse_office_document(file_path: str, mime_type: str) -> str:
    """
    Intelligently route to the appropriate office document parser.
    Returns extracted plain text or markdown string.
    """
    path_obj = Path(file_path)
    extension = path_obj.suffix.lower()
    
    try:
        if extension == ".pdf" or mime_type == "application/pdf":
            return parse_pdf(path_obj)
        elif extension == ".docx" or "wordprocessingml" in mime_type:
            return parse_docx(path_obj)
        elif extension == ".pptx" or "presentationml" in mime_type:
            return parse_pptx(path_obj)
        elif extension == ".xlsx" or "spreadsheetml" in mime_type:
            return parse_xlsx(path_obj)
        else:
            raise ValueError(f"Unsupported office document type: {extension} / {mime_type}")
    except Exception as e:
        print(f"[Advanced Skills Error] Failed to parse {file_path}: {e}")
        return ""
